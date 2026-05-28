import asyncio
from pathlib import Path

from app.core.logging import get_logger
from app.extraction.base import BaseExtractor, ExtractionResult, TableData

logger = get_logger(__name__)


class PptxExtractor(BaseExtractor):
    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx"]

    async def extract(self, file_path: Path) -> ExtractionResult:
        return await asyncio.to_thread(self._extract_sync, file_path)

    def _extract_sync(self, file_path: Path) -> ExtractionResult:
        from pptx import Presentation

        result = ExtractionResult()
        prs = Presentation(str(file_path))
        all_text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_text_parts.append(text)

                if shape.has_table:
                    table = shape.table
                    headers = [cell.text.strip() for cell in table.rows[0].cells]
                    rows = []
                    for row in table.rows[1:]:
                        rows.append([cell.text.strip() for cell in row.cells])
                    result.tables.append(
                        TableData(headers=headers, rows=rows, page_number=slide_num)
                    )

            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                all_text_parts.append(f"Slide {slide_num}:\n{slide_text}")
                result.pages.append({"page": slide_num, "text": slide_text})

            # Extract notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    all_text_parts.append(f"Slide {slide_num} Notes:\n{notes}")

        result.text = "\n\n".join(all_text_parts)
        result.metadata = {
            "slide_count": len(prs.slides),
            "word_count": len(result.text.split()),
            "table_count": len(result.tables),
        }
        return result
