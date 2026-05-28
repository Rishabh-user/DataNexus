"""PPT generation orchestrator — theme-aware."""
from pathlib import Path

from pptx import Presentation

from app.core.logging import get_logger
from app.ppt.charts import add_chart_slide
from app.ppt.content_formatter import format_bullet_points, format_table_data
from app.ppt.styles import DEFAULT_THEME_ID, THEMES
from app.ppt.templates import (
    add_content_slide,
    add_summary_slide,
    add_table_slide,
    add_title_slide,
)

logger = get_logger(__name__)


def generate_ppt(
    slide_data: dict,
    output_path: str,
    include_charts: bool = True,
    theme_id: str = DEFAULT_THEME_ID,
) -> str:
    theme = THEMES.get(theme_id) or THEMES[DEFAULT_THEME_ID]

    prs = Presentation()
    prs.slide_width  = 12192000
    prs.slide_height = 6858000

    title  = slide_data.get("title", "Report")
    slides = slide_data.get("slides", [])

    if not isinstance(slides, list):
        slides = []
    if not slides:
        slides = [{"type": "title", "title": title, "bullet_points": ["Generated Report"]}]

    for slide_info in slides:
        slide_type  = slide_info.get("type", "content")
        slide_title = slide_info.get("title", "")
        notes       = slide_info.get("notes", "")

        try:
            if slide_type == "title":
                subtitle = (slide_info.get("bullet_points") or [""])[0]
                add_title_slide(prs, slide_title, subtitle, theme)

            elif slide_type == "chart" and include_charts:
                chart_data = slide_info.get("chart_data")
                if chart_data:
                    add_chart_slide(prs, slide_title, chart_data, theme, notes)
                else:
                    add_content_slide(prs, slide_title,
                        format_bullet_points(slide_info.get("bullet_points", [])),
                        theme, notes)

            elif slide_type == "table":
                table_data = slide_info.get("table_data")
                if table_data:
                    headers, rows = format_table_data(table_data)
                    add_table_slide(prs, slide_title, headers, rows, theme)
                else:
                    add_content_slide(prs, slide_title,
                        format_bullet_points(slide_info.get("bullet_points", [])),
                        theme, notes)

            elif slide_type == "summary":
                add_summary_slide(prs, slide_title,
                    format_bullet_points(slide_info.get("bullet_points", [])), theme)

            else:
                add_content_slide(prs, slide_title,
                    format_bullet_points(slide_info.get("bullet_points", [])),
                    theme, notes)

        except Exception as exc:
            import traceback
            logger.warning(
                "Failed to create slide '%s': %s\n%s",
                slide_title, exc, traceback.format_exc()
            )

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    logger.info("PPT saved (%s): %s — %d slides", theme.name, output_path, len(prs.slides))
    return output_path
