"""Slide builder functions — professionally designed, theme-aware."""
from __future__ import annotations

from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.ppt.styles import (
    FONT_BODY, FONT_FAMILY, FONT_HEADING, FONT_SMALL,
    FONT_SUBTITLE, FONT_TABLE, FONT_TITLE, Theme,
)

Rgb = tuple[int, int, int]

# ─────────────────────────────────────────────────────────────────────────────
# Primitive helpers
# ─────────────────────────────────────────────────────────────────────────────

def _rgb(t: Rgb) -> RGBColor:
    return RGBColor(*t)


def _blend(c1: Rgb, c2: Rgb, t: float) -> Rgb:
    """Blend c1 towards c2 by factor t  (0 = c1, 1 = c2)."""
    return tuple(int(c1[i] * (1 - t) + c2[i] * t) for i in range(3))


def _is_dark(color: Rgb) -> bool:
    r, g, b = color
    return (0.299 * r + 0.587 * g + 0.114 * b) < 128


def _set_bg(slide, color: Rgb) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide, left: float, top: float, width: float, height: float, fill: Rgb):
    s = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    s.line.fill.background()
    return s


def _oval(slide, left: float, top: float, width: float, height: float, fill: Rgb):
    s = slide.shapes.add_shape(
        9,  # OVAL
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = _rgb(fill)
    s.line.fill.background()
    return s


def _txbox(slide, left: float, top: float, width: float, height: float, wrap: bool = True):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = tb.text_frame
    tf.word_wrap = wrap
    return tf


def _run(para, text: str, size: float, *,
         bold: bool = False, italic: bool = False, color: Rgb | None = None):
    r = para.add_run()
    r.text = text
    r.font.name = FONT_FAMILY
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    if color:
        r.font.color.rgb = _rgb(color)
    return r


def _add_picture_safe(slide, buf: BytesIO | None,
                      left: float, top: float, width: float, height: float) -> None:
    if buf is None:
        return
    try:
        slide.shapes.add_picture(buf, Inches(left), Inches(top), Inches(width), Inches(height))
    except Exception:
        pass


# ─────────────────────────────────────────────────────────────────────────────
# Shared layout elements
# ─────────────────────────────────────────────────────────────────────────────

def _slide_header(slide, title: str, theme: Theme) -> None:
    # Full-height left stripe
    _rect(slide, 0, 0, 0.13, 7.5, theme.accent)
    # Header band
    _rect(slide, 0.13, 0, 13.22, 1.30, theme.title_bg)
    # Accent underline
    _rect(slide, 0.13, 1.30, 13.22, 0.055, theme.accent)
    # Accent2 type-indicator bar
    _rect(slide, 0.26, 0.22, 0.09, 0.86, theme.accent2)
    # Title text
    tf = _txbox(slide, 0.52, 0.20, 12.55, 0.94)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _run(p, title, FONT_HEADING, bold=True, color=theme.title_text)


def _slide_footer(slide, theme: Theme, label: str = "DataNexus AI Report") -> None:
    footer_bg = _blend(theme.title_bg, (0, 0, 0), 0.18)
    _rect(slide, 0, 7.07, 13.333, 0.43, footer_bg)
    fg = theme.accent if _is_dark(footer_bg) else theme.heading_text
    # Left label
    tf1 = _txbox(slide, 0.28, 7.10, 8.5, 0.34)
    p1 = tf1.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    _run(p1, label, FONT_SMALL - 1, color=fg)
    # Right branding
    tf2 = _txbox(slide, 9.2, 7.10, 3.9, 0.34)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    _run(p2, "Powered by DataNexus AI", FONT_SMALL - 1, color=theme.accent)


# ─────────────────────────────────────────────────────────────────────────────
# TITLE SLIDE
# ─────────────────────────────────────────────────────────────────────────────

def add_title_slide(prs: Presentation, title: str, subtitle: str, theme: Theme) -> None:
    from app.ppt.image_gen import make_title_deco

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.title_bg)

    # Decorative PNG — right side, drawn first (behind everything)
    deco = make_title_deco(theme, width_px=700, height_px=520)
    _add_picture_safe(slide, deco, 6.8, -0.25, 6.5, 4.85)

    # Left accent stripe
    _rect(slide, 0, 0, 0.18, 7.5, theme.accent)

    # Top branding bar
    top_bar_bg = _blend(theme.title_bg, (0, 0, 0), 0.30)
    _rect(slide, 0.18, 0, 13.15, 0.54, top_bar_bg)
    brand_fg = theme.accent if _is_dark(top_bar_bg) else theme.heading_text
    tf_brand = _txbox(slide, 0.38, 0.09, 5.5, 0.36)
    p_brand = tf_brand.paragraphs[0]
    p_brand.alignment = PP_ALIGN.LEFT
    _run(p_brand, "DataNexus AI", FONT_SMALL, bold=True, color=brand_fg)
    tf_date = _txbox(slide, 9.6, 0.10, 3.65, 0.32)
    p_date = tf_date.paragraphs[0]
    p_date.alignment = PP_ALIGN.RIGHT
    _run(p_date, datetime.now().strftime("%B %d, %Y"), FONT_SMALL - 1, color=brand_fg)

    # Main title
    tf_title = _txbox(slide, 0.45, 1.45, 9.0, 2.75)
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    p_title.space_after = Pt(6)
    _run(p_title, title, FONT_TITLE, bold=True, color=theme.title_text)

    # Accent divider
    _rect(slide, 0.45, 4.42, 4.8, 0.08, theme.accent)

    # Subtitle
    if subtitle:
        tf_sub = _txbox(slide, 0.45, 4.60, 8.5, 1.3)
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.alignment = PP_ALIGN.LEFT
        sub_fg = theme.accent2 if _is_dark(theme.title_bg) else theme.heading_text
        _run(p_sub, subtitle, FONT_SUBTITLE, color=sub_fg)

    # Bottom bar
    _rect(slide, 0.18, 6.85, 13.15, 0.65, theme.accent2)
    tf_bot = _txbox(slide, 0.42, 6.90, 12.5, 0.50)
    p_bot = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    bot_fg = (255, 255, 255) if _is_dark(theme.accent2) else theme.body_text
    _run(p_bot, "AI-Powered Data Intelligence  •  DataNexus", FONT_SMALL, color=bot_fg)


# ─────────────────────────────────────────────────────────────────────────────
# CONTENT SLIDE
# ─────────────────────────────────────────────────────────────────────────────

def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    theme: Theme,
    notes: str = "",
) -> None:
    from app.ppt.image_gen import make_corner_deco

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.slide_bg)

    # Corner decoration (bottom-right, drawn before header)
    deco = make_corner_deco(theme, size_px=260)
    _add_picture_safe(slide, deco, 11.1, 4.8, 2.1, 2.1)

    # Standard header
    _slide_header(slide, title, theme)

    # Numbered bullet items
    items      = (bullets or ["No content available."])[:8]
    n          = len(items)
    y0         = 1.45
    y_end      = 7.07
    avail      = y_end - y0
    per_h      = min(avail / n, 0.82)
    badge_size = 0.33
    badge_fg   = (255, 255, 255) if _is_dark(theme.accent) else (20, 20, 40)

    for i, item in enumerate(items):
        y  = y0 + i * per_h
        by = y + (per_h - badge_size) / 2

        # Numbered oval badge
        _oval(slide, 0.22, by, badge_size, badge_size, theme.accent)
        tf_num = _txbox(slide, 0.22, by + 0.01, badge_size, badge_size - 0.04)
        p_num  = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        _run(p_num, str(i + 1), FONT_SMALL - 1, bold=True, color=badge_fg)

        # Thin separator rule (except after last item)
        if i < n - 1:
            rule_col = _blend(theme.accent, theme.slide_bg, 0.82)
            _rect(slide, 0.70, y + per_h - 0.025, 12.2, 0.016, rule_col)

        # Bullet text
        tf_txt = _txbox(slide, 0.70, y + 0.04, 12.35, per_h - 0.08)
        tf_txt.word_wrap = True
        p_txt  = tf_txt.paragraphs[0]
        p_txt.alignment = PP_ALIGN.LEFT
        _run(p_txt, str(item).strip(), FONT_BODY, color=theme.body_text)

    _slide_footer(slide, theme)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ─────────────────────────────────────────────────────────────────────────────
# TABLE SLIDE
# ─────────────────────────────────────────────────────────────────────────────

def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list],
    theme: Theme,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.slide_bg)

    _slide_header(slide, title, theme)

    if not headers or not rows:
        tf = _txbox(slide, 0.5, 1.7, 12.0, 1.0)
        _run(tf.paragraphs[0], "No table data available.", FONT_BODY, color=theme.body_text)
        _slide_footer(slide, theme)
        return

    n_rows    = min(len(rows), 16)
    n_cols    = len(headers)
    show_rows = rows[:n_rows]

    tbl_x = 0.28
    tbl_y = 1.42
    tbl_w = 12.77
    row_h = 0.37
    tbl_h = min(row_h * (n_rows + 1), 5.40)

    tbl_shape = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(tbl_x), Inches(tbl_y),
        Inches(tbl_w), Inches(tbl_h),
    )
    tbl = tbl_shape.table

    for ci in range(n_cols):
        tbl.columns[ci].width = int(12192000 // n_cols)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(theme.table_header_bg)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(hdr)[:28]
        r.font.name = FONT_FAMILY
        r.font.size = Pt(FONT_TABLE + 1)
        r.font.bold = True
        r.font.color.rgb = _rgb(theme.table_header_fg)

    # Data rows
    for ri, row in enumerate(show_rows):
        is_alt = ri % 2 == 1
        row_bg = theme.table_alt_row if is_alt else theme.slide_bg
        for ci in range(n_cols):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(row_bg)
            val = row[ci] if ci < len(row) else ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)[:42] if val is not None else ""
            r.font.name = FONT_FAMILY
            r.font.size = Pt(FONT_TABLE)
            r.font.color.rgb = _rgb(theme.body_text)

    # Clipped-rows caption
    if len(rows) > n_rows:
        cap_y = tbl_y + tbl_h + 0.06
        tf_cap = _txbox(slide, tbl_x, cap_y, tbl_w, 0.28)
        p_cap  = tf_cap.paragraphs[0]
        p_cap.alignment = PP_ALIGN.RIGHT
        _run(p_cap, f"Showing {n_rows} of {len(rows)} records",
             FONT_SMALL - 1, italic=True, color=theme.body_text)

    _slide_footer(slide, theme)


# ─────────────────────────────────────────────────────────────────────────────
# SUMMARY SLIDE
# ─────────────────────────────────────────────────────────────────────────────

def add_summary_slide(
    prs: Presentation,
    title: str,
    points: list[str],
    theme: Theme,
) -> None:
    from app.ppt.image_gen import make_summary_deco

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, theme.title_bg)

    # Background depth ovals
    _oval(slide,  8.6, -1.4, 5.8, 5.8, _blend(theme.accent,  theme.title_bg, 0.60))
    _oval(slide, 10.2,  3.8, 4.2, 4.2, _blend(theme.accent2, theme.title_bg, 0.65))
    _oval(slide, -1.8,  4.8, 4.5, 4.5, _blend(theme.accent,  theme.title_bg, 0.70))

    # Decorative PNG (right side)
    deco = make_summary_deco(theme, width_px=560, height_px=440)
    _add_picture_safe(slide, deco, 9.0, 0.6, 4.2, 3.3)

    # Left stripe
    _rect(slide, 0, 0, 0.18, 7.5, theme.accent)

    # Top accent bar
    _rect(slide, 0.18, 0, 13.15, 0.17, theme.accent2)

    # Title panel
    title_panel_bg = _blend(theme.title_bg, (0, 0, 0), 0.20)
    _rect(slide, 0.18, 0.17, 13.15, 1.10, title_panel_bg)
    tf_title = _txbox(slide, 0.48, 0.24, 11.5, 0.96)
    p_title  = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    _run(p_title, title or "Key Takeaways", FONT_HEADING + 4, bold=True, color=theme.title_text)

    # Divider
    _rect(slide, 0.48, 1.33, 11.8, 0.06, theme.accent)

    # Numbered items
    items      = (points or ["Summary not available."])[:7]
    n          = len(items)
    y0         = 1.50
    y_end      = 6.85
    per_h      = min((y_end - y0) / n, 0.84)
    badge_size = 0.34
    badge_fg   = (255, 255, 255) if _is_dark(theme.accent) else (20, 20, 40)

    for i, pt in enumerate(items):
        y  = y0 + i * per_h
        by = y + (per_h - badge_size) / 2

        _oval(slide, 0.30, by, badge_size, badge_size, theme.accent)
        tf_num = _txbox(slide, 0.30, by + 0.01, badge_size, badge_size - 0.04)
        p_num  = tf_num.paragraphs[0]
        p_num.alignment = PP_ALIGN.CENTER
        _run(p_num, str(i + 1), FONT_SMALL - 1, bold=True, color=badge_fg)

        tf_txt = _txbox(slide, 0.80, y + 0.04, 8.4, per_h - 0.06)
        tf_txt.word_wrap = True
        p_txt  = tf_txt.paragraphs[0]
        p_txt.alignment = PP_ALIGN.LEFT
        _run(p_txt, str(pt).strip(), FONT_BODY, color=theme.title_text)

    # Bottom bar
    _rect(slide, 0.18, 6.85, 13.15, 0.65, theme.accent2)
    tf_bot = _txbox(slide, 0.42, 6.90, 12.5, 0.52)
    p_bot  = tf_bot.paragraphs[0]
    p_bot.alignment = PP_ALIGN.LEFT
    bot_fg = (255, 255, 255) if _is_dark(theme.accent2) else theme.body_text
    _run(p_bot, "DataNexus AI  •  AI-Powered Data Intelligence", FONT_SMALL, color=bot_fg)
