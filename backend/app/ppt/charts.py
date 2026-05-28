"""Chart slide builder — theme-aware with professional header/footer."""
from __future__ import annotations

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

from app.ppt.styles import FONT_FAMILY, FONT_SMALL, Theme


def _rgb(t):
    return RGBColor(*t)


def _safe_num(v) -> float:
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def add_chart_slide(
    prs: Presentation,
    title: str,
    chart_data: dict,
    theme: Theme,
    notes: str = "",
) -> None:
    from app.ppt.templates import _slide_header, _slide_footer, _txbox, _run

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = _rgb(theme.slide_bg)

    # Standard header + footer
    _slide_header(slide, title, theme)

    chart_type_str = (chart_data.get("type") or "bar").lower()
    chart_type_map = {
        "bar":    XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line":   XL_CHART_TYPE.LINE_MARKERS,
        "pie":    XL_CHART_TYPE.PIE,
    }
    xl_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    labels   = chart_data.get("labels", [])
    datasets = chart_data.get("datasets", [])

    if not labels or not datasets:
        tf = _txbox(slide, 0.5, 1.7, 12.0, 1.0)
        _run(tf.paragraphs[0], "No chart data available.", 14, color=theme.body_text)
        _slide_footer(slide, theme)
        return

    cd = ChartData()
    cd.categories = labels

    for i, ds in enumerate(datasets):
        if isinstance(ds, dict):
            label    = ds.get("label", f"Series {i + 1}")
            raw_vals = ds.get("values", [])
        else:
            label    = f"Series {i + 1}"
            raw_vals = ds if isinstance(ds, list) else []

        if raw_vals and isinstance(raw_vals[0], list):
            for j, sub in enumerate(raw_vals):
                cd.add_series(f"Series {j + 1}", [_safe_num(v) for v in sub])
        else:
            cd.add_series(label, [_safe_num(v) for v in raw_vals])

    # Chart area — sits between header (1.38") and footer (7.07")
    chart_shape = slide.shapes.add_chart(
        xl_type,
        Inches(0.30), Inches(1.42),
        Inches(12.75), Inches(5.50),
        cd,
    )
    chart = chart_shape.chart

    # Apply theme colours to series
    for i, series in enumerate(chart.series):
        color = theme.chart_colors[i % len(theme.chart_colors)]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = _rgb(color)
        if chart_type_str == "line":
            series.format.line.color.rgb = _rgb(color)
            series.format.line.width = Pt(2.25)

    # Legend styling
    if chart.has_legend:
        chart.legend.font.size = Pt(FONT_SMALL)
        chart.legend.font.name = FONT_FAMILY

    # Plot area background = slide bg
    try:
        chart.plot_area.format.fill.solid()
        chart.plot_area.format.fill.fore_color.rgb = _rgb(theme.slide_bg)
    except Exception:
        pass

    _slide_footer(slide, theme)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
